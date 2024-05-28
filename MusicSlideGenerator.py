import os
from tkinter import Tk, Label, Entry, Button, Listbox, SINGLE, MULTIPLE, END, filedialog, messagebox
from pptx import Presentation
from fuzzywuzzy import process

class MusicSlidesGenerator:
    def __init__(self):
        self.root = Tk()
        self.root.title("Church Music Slides Generator")
        self.create_widgets()
        self.drag_data = {"x": 0, "y": 0, "item": None}

    def create_widgets(self):
        # Directory selection
        Label(self.root, text="Select Directory:").grid(row=0, column=0, padx=10, pady=10)
        self.directory_entry = Entry(self.root, width=50)
        self.directory_entry.grid(row=0, column=1, padx=10, pady=10)
        Button(self.root, text="Browse", command=self.browse_directory).grid(row=0, column=2, padx=10, pady=10)

        # Fuzzy search input
        Label(self.root, text="Search Files:").grid(row=1, column=0, padx=10, pady=10)
        self.search_entry = Entry(self.root, width=50)
        self.search_entry.grid(row=1, column=1, padx=10, pady=10)
        Button(self.root, text="Search", command=self.search_files).grid(row=1, column=2, padx=10, pady=10)
        self.search_entry.bind('<Return>', self.on_enter_key)

        # Listbox for search results
        self.results_listbox = Listbox(self.root, selectmode=SINGLE, width=50, height=10)
        self.results_listbox.grid(row=2, column=1, padx=10, pady=10)
        Button(self.root, text="Add Selected Files", command=self.add_selected_files).grid(row=2, column=2, padx=10, pady=10)
        self.results_listbox.bind('<Double-1>', self.on_double_click)
        
        # # Song titles input
        # Label(self.root, text="Song Title:").grid(row=3, column=0, padx=10, pady=10)
        # self.song_entry = Entry(self.root, width=50)
        # self.song_entry.grid(row=1, column=1, padx=10, pady=10)
        # Button(self.root, text="Add Song", command=self.add_song).grid(row=1, column=2, padx=10, pady=10)

        # Listbox for song titles
        self.song_listbox = Listbox(self.root, selectmode=SINGLE, width=50, height=10)
        self.song_listbox.grid(row=4, column=1, padx=10, pady=10)
        Button(self.root, text="Remove Selected Songs", command=self.remove_selected_songs).grid(row=4, column=2, padx=10, pady=10)
        self.song_listbox.bind("<Button-1>", self.on_click)
        self.song_listbox.bind("<B1-Motion>", self.on_drag)
        self.song_listbox.bind("<ButtonRelease-1>", self.on_drop)

        # Generate presentation button
        Button(self.root, text="Generate Presentation", command=self.generate_presentation).grid(row=5, column=1, padx=10, pady=10)

    def browse_directory(self):
        directory = filedialog.askdirectory()
        if directory:
            self.directory_entry.delete(0, END)
            self.directory_entry.insert(0, directory)

    # def add_song(self):
    #     song = self.song_entry.get().strip()
    #     if song:
    #         self.song_listbox.insert(END, song)
    #         self.song_entry.delete(0, END)

    def search_files(self):
        directory = self.directory_entry.get().strip()
        if not directory:
            messagebox.showwarning("Missing Directory", "Please select a directory.")
            return

        search_term = self.search_entry.get().strip()
        if not search_term:
            messagebox.showwarning("Missing Search Term", "Please enter a search term.")
            return

        try:
            files = [f for f in os.listdir(directory) if f.endswith('.pptx')]
            matches = process.extract(search_term, files, limit=10)
            self.results_listbox.delete(0, END)
            for match in matches:
                self.results_listbox.insert(END, match[0])
        except Exception as e:
            messagebox.showerror("Error", f"An error occurred while searching: {e}")

    def on_enter_key(self, event):
        self.search_files()
    
    def on_double_click(self, event):
        selected_files = self.results_listbox.curselection()
        for index in selected_files:
            file_name = self.results_listbox.get(index)
            self.add_to_listbox(self.song_listbox, file_name)

    def add_selected_files(self):
        selected_files = self.results_listbox.curselection()
        for index in selected_files:
            file_name = self.results_listbox.get(index)
            self.add_to_listbox(self.song_listbox, file_name)
            
    def add_to_listbox(self, listbox, item):
        if item in listbox.get(0, END):
            if not messagebox.askyesno("Duplicate Entry", f"'{item}' is already in the list. Do you want to add it again?"):
                return
        listbox.insert(END, item)
        
    def on_click(self, event):
        widget = event.widget
        self.drag_data["item"] = widget.nearest(event.y)
        self.drag_data["y"] = event.y

    def on_drag(self, event):
        widget = event.widget
        new_index = widget.nearest(event.y)
        if new_index != self.drag_data["item"]:
            item_text = widget.get(self.drag_data["item"])
            widget.delete(self.drag_data["item"])
            widget.insert(new_index, item_text)
            self.drag_data["item"] = new_index
        self.drag_data["y"] = event.y

    def on_drop(self, event):
        self.drag_data["item"] = None
            
    def remove_selected_songs(self):
        selected_songs = self.song_listbox.curselection()
        for index in reversed(selected_songs):
            self.song_listbox.delete(index)

    def generate_presentation(self):
        directory = self.directory_entry.get().strip()
        song_titles = list(self.song_listbox.get(0, END))
        output_file = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint files", "*.pptx")])

        if not directory or not song_titles or not output_file:
            messagebox.showwarning("Missing Information", "Please provide all required inputs.")
            return

        try:
            self.concatenate_presentations(directory, song_titles, output_file)
            messagebox.showinfo("Success", "Presentation created successfully!")
        except Exception as e:
            messagebox.showerror("Error", f"An error occurred: {e}")

    @staticmethod
    def get_slide_title(slide):
        for shape in slide.shapes:
            if shape.has_text_frame:
                return shape.text_frame.text.strip()
        return None

    @staticmethod
    def add_title_slide(prs, title):
        slide_layout = prs.slide_layouts[0]  # Use the title slide layout
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title

    def concatenate_presentations(self, directory, song_titles, output_file):
        presentations = [os.path.join(directory, f) for f in os.listdir(directory) if f.endswith('.pptx')]
        combined_prs = Presentation()

        for song_title in song_titles:
            self.add_title_slide(combined_prs, song_title)
            for presentation in presentations:
                prs = Presentation(presentation)
                for slide in prs.slides:
                    slide_title = self.get_slide_title(slide)
                    if slide_title and slide_title.lower() == song_title.lower():
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                blank_slide_layout = combined_prs.slide_layouts[5]  # Use the blank slide layout
                                new_slide = combined_prs.slides.add_slide(blank_slide_layout)
                                for shape in slide.shapes:
                                    new_shape = new_slide.shapes.add_shape(shape.auto_shape_type, shape.left, shape.top, shape.width, shape.height)
                                    new_shape.text = shape.text

        combined_prs.save(output_file)